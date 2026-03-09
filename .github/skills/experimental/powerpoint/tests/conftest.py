"""Shared fixtures for PowerPoint skill tests."""

import struct
import zlib

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


@pytest.fixture()
def blank_presentation():
    """Fresh Presentation with standard widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


@pytest.fixture()
def blank_slide(blank_presentation):
    """Blank slide added to a fresh presentation."""
    layout = blank_presentation.slide_layouts[6]  # Blank layout
    return blank_presentation.slides.add_slide(layout)


@pytest.fixture()
def sample_textbox(blank_slide):
    """Slide with a textbox containing known text and formatting."""
    txBox = blank_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.text = "Sample Text"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True
    return txBox


@pytest.fixture()
def sample_shape(blank_slide):
    """Slide with a rectangle shape having fill and text."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = blank_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2),
        Inches(2),
        Inches(3),
        Inches(2),
    )
    shape.text = "Shape Text"
    shape.fill.solid()
    from pptx.dml.color import RGBColor

    shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
    return shape


def _minimal_png_bytes() -> bytes:
    """Create a minimal valid 1x1 red PNG in memory."""

    def _chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw_row = b"\x00\xff\x00\x00"  # filter byte + RGB
    idat_data = zlib.compress(raw_row)
    return (
        signature
        + _chunk(b"IHDR", ihdr_data)
        + _chunk(b"IDAT", idat_data)
        + _chunk(b"IEND", b"")
    )


@pytest.fixture()
def sample_image_path(tmp_path):
    """Minimal valid PNG file at a temporary path."""
    img = tmp_path / "test.png"
    img.write_bytes(_minimal_png_bytes())
    return img
